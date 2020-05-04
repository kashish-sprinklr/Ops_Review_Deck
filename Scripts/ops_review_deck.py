from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pathlib import Path
import os
import time
import pandas as pd

legend_dict = {0: {"width": 10, "cell_color": (251, 251, 251)},
               1: {"width": Cm(0.2), "cell_color": (191, 191, 191)},
               2: {"width": Cm(2.5), "align": PP_ALIGN.CENTER, "cell_color": (251, 251, 251), "text": "Average EHI",
                   "font_size": 10, "bold": False, "font_color": (123, 123, 123)},
               3: {"width": 10, "cell_color": (251, 251, 251)},
               4: {"align": PP_ALIGN.CENTER, "rotation": True, "cell_color": (251, 251, 251), "text": "|",
                   "font_size": 25, "bold": True, "font_color": (155, 187, 89)},
               5: {"align": PP_ALIGN.CENTER, "width": Inches(1.55), "cell_color": (251, 251, 251),
                   "text": "#High EHI(8,9,10)", "font_size": 10, "bold": False, "font_color": (123, 123, 123)},
               6: {"align": PP_ALIGN.CENTER, "rotation": True, "cell_color": (251, 251, 251), "text": "|",
                   "font_size": 25, "bold": True, "font_color": (247, 150, 70)},
               7: {"align": PP_ALIGN.CENTER, "width": Inches(1.7), "cell_color": (251, 251, 251),
                   "text": "#Medium EHI(5,6,7)", "font_size": 10, "bold": False, "font_color": (123, 123, 123)},
               8: {"align": PP_ALIGN.CENTER, "rotation": True, "cell_color": (251, 251, 251), "text": "|",
                   "font_size": 25, "bold": True, "font_color": (193, 81, 78)},
               9: {"align": PP_ALIGN.CENTER, "width": Inches(1.7), "cell_color": (251, 251, 251),
                   "text": "#Low EHI(1,2,3,4)", "font_size": 10, "bold": False, "font_color": (123, 123, 123)}}


def get_dataframe_dict(path, file_names_list):
    df_dict = dict()
    for file in file_names_list:
        ref_file = "_".join(file.lower().split(" "))
        df = pd.read_excel(os.path.join(path, file + ".xlsx"))
        df = df.drop(columns=["Unnamed: 0"])
        df = df.rename(columns={"{}_headcount".format(ref_file): "headcount"})
        df_dict[file] = df
    return df_dict


def filter_dataframe_dict_on_uid(_id, df_dict):
    df_fdict = dict()
    for key in df_dict.keys():
        try:
            df = df_dict[key]
            df = dict(df[df["associate_id"] == _id].iloc[-1])
            df_fdict[key] = df
        except:
            pass
    return df_fdict


def get_title(df_fdict):
    for key in df_fdict.keys():
        title = df_fdict[key]["slt_name"]
        return title


def get_json_data_for_table_1(dbs, df_fdict):
    table_dict = {"KPI": ["Headcount", "#Regrettable Voluntary Attrition", "Regrettable Voluntary Attrition",
                          "Employee Delight Assurance Program (EDAP)", "#Promotion",
                          "Female Leadership Representation"]}
    for db_name in dbs:
        df = df_fdict[db_name]
        hc = df["headcount"]
        nrve = df["num_regrettable_voluntary_exits"]
        rvar = df["regrettable_voluntary_attrition_rate"]
        es = df["e_score"]
        np = df["num_promotions"]
        plf = df["perc_leaders_female"]
        values = []
        for i, val in enumerate([hc, nrve, rvar, es, np, plf]):
            val = str(val) + "%" if i in [2, 5] else str(val)
            values.append(val)
        table_dict["Actual\n {}".format(db_name)] = values
    return table_dict


def get_json_data_for_table_2(df_fdict):
    table_dict = {"Categories": ["Total Annualized Attrition", "Voluntary\nNon Regrettable", "Regrettable",
                                 "Involuntary"]}
    for db_name in df_fdict.keys():
        df = df_fdict[db_name]
        base = df["headcount"]
        aar = df["annualized_attrition_rate"]  # to be changed already in percentage
        nnrve = df["num_non_regrettable_voluntary_exits"]
        nrve = df["num_regrettable_voluntary_exits"]
        nie = df["num_involuntary_exits"]
        table_dict[db_name] = [str(x) + "%" for x in [aar]] + [str(round((val / base) * 100, 2)) +
                                                               "%" for val in [nnrve, nrve, nie]]
    return table_dict


def read_ehi_data(file_path):
    df = pd.read_excel(file_path)
    df = df.dropna(how="all", axis=0)
    df.columns = df.iloc[0].fillna('Category')
    df.columns = ["Category"] + [int(x) for x in df.columns if x not in ["Category"]]
    df = df.iloc[1:]
    return df


def sort_ehi_category(y_axis, values):
    new_y_axis, new_values = [0] * 3, [0] * 3
    for i, cat in enumerate(["High", "Medium", "Low"]):
        new_y_axis[i] = y_axis[y_axis.index(cat)]
        new_values[i] = values[y_axis.index(cat)]
    return new_y_axis, new_values


def get_dict_attr_from_dataframe_for_ehi(df):
    cycle_dict = {1: "1 (Jul'18)", 2: "2 (Sep'18)", 3: "3 (Dec'18)", 4: "4 (Feb'19)", 5: "5 (Apr'19)",
                  6: "6 (Jun'19)", 7: "7 (Sep'19)", 8: "8 (Dec'19)", 9: "9 (Feb'20)", 10: "10 (Jun'20)",
                  11: "11 (Sep'20)", 12: "12 (Dec'20)", 13: "13 (Feb'21)", 14: "14 (Jun'21)"}
    df_cycles = sorted([int(x) for x in df.columns if x not in ["Category"]])
    x_axis = [cycle_dict[x] for x in df_cycles]
    y_axis_1 = ["EHI"]
    y_axis_2 = [x for x in df["Category"].unique() if x not in ["E-Score"]]
    values_1 = [str(x) for x in df[df["Category"] == "E-Score"].values[0][1:]]
    values_2 = []
    for cat in y_axis_2:
        values_2.append(list(df[df["Category"] == cat].values[0][1:]))
    if len(y_axis_2) == 1:
        y_axis_2 = y_axis_2 + [x for x in ["High", "Medium", "Low"] if x not in y_axis_2]
        zero_list = [[0] * len(x_axis) for _ in range(len(y_axis_2) - len(values_2))]
        values_2 = values_2 + zero_list
    y_axis_2, values_2 = sort_ehi_category(y_axis_2, values_2)
    data_labels = []
    for lst in values_2:
        if sum(lst) == 0:
            data_labels.append(False)
        else:
            data_labels.append(True)
    return x_axis, y_axis_1, [values_1], y_axis_2, values_2, data_labels


def fill_text_color(font, rgb):
    fill = font.fill
    fill.patterned()
    fillcolor = fill.fore_color
    fillcolor.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def add_text_box(slide, box_dim, text, pt, bold=False, style="Proxima Nova", align_center=True, rotation=0,
                 rgb=(0, 0, 0)):
    text_box = slide.shapes.add_textbox(box_dim[0], box_dim[1], box_dim[2], box_dim[3])
    text_box.rotation = rotation
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    if align_center:
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.name = style
    run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    return run.font


def format_cell_text(cell, text_ptr, data, pt, bold=False, style="Proxima Nova", align=None):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_ptr.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = data
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.name = style
    return run.font


def fill_table_cell_color(cell, rgb):
    fill = cell.fill
    fill.solid()
    fillcolor = fill.fore_color
    fillcolor.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def fill_chart_series_color(points, rgb):
    fill = points.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def insert_table_content(table, table_content, font_size):
    for i, (key, values) in enumerate(table_content.items()):
        cell = table.cell(0, i)
        font = format_cell_text(cell, cell.text_frame, key, font_size, True, align=PP_ALIGN.CENTER)
        fill_text_color(font, [255, 255, 255])
        fill_table_cell_color(cell, [68, 85, 105])
        for j, data in enumerate(values):
            cell = table.cell(j + 1, i)
            if i == 0:
                format_cell_text(cell, cell.text_frame, data, font_size)
            else:
                format_cell_text(cell, cell.text_frame, data, font_size, align=PP_ALIGN.CENTER)
            if j % 2 == 0:
                fill_table_cell_color(cell, [231, 230, 230])
            else:
                fill_table_cell_color(cell, [252, 252, 252])


def create_table_box(slide, table_dim, font_size, table_content, fix_column=None):
    columns = list(table_content.keys())
    shape = slide.shapes.add_table(len(table_content[columns[0]]) + 1, len(columns), table_dim[0], table_dim[1],
                                   table_dim[2], table_dim[3])
    table = shape.table
    insert_table_content(table, table_content, font_size)
    if fix_column:
        table.columns[0].width = Inches(fix_column)


def reformat_dict_attr_for_attrition_chart(df_fdict):
    objs = ["Involuntary", "Voluntary Regrettable", "Voluntary Non Regrettable"]
    nie, nrve, nnrve = [], [], []
    for key in df_fdict.keys():
        temp_db = df_fdict[key]
        nie.append(temp_db["num_involuntary_exits"])
        nrve.append(temp_db["num_regrettable_voluntary_exits"])
        nnrve.append(temp_db["num_non_regrettable_voluntary_exits"])
    keys = list(df_fdict.keys())
    values = [nie, nrve, nnrve]
    new_objs, new_values = [], []
    zero_objs, zero_values = [], []
    data_label_flag_1, data_label_flag_2 = [], []
    for val_list, obj in zip(values, objs):
        if sum(val_list) != 0:
            new_objs.append(obj)
            new_values.append(val_list)
            data_label_flag_1.append(True)
        else:
            zero_objs.append(obj)
            zero_values.append(val_list)
            data_label_flag_2.append(False)
    objs = zero_objs + new_objs
    values = zero_values + new_values
    data_label_flag = data_label_flag_2 + data_label_flag_1
    return keys, objs, values, data_label_flag


def format_chart_1_properties(chart, y_axis, label_flag=None):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(13)
    chart.value_axis.has_major_gridlines = False
    tick_label = chart.value_axis.tick_labels
    tick_label.font.size = Pt(14)
    chart.category_axis.has_major_gridlines = False
    tick_label = chart.category_axis.tick_labels
    tick_label.font.size = Pt(14)
    plot = chart.plots[0]
    series_color = [(218, 227, 243), (0, 176, 240), (166, 166, 166)]
    for i in range(len(y_axis)):
        points = plot.series[i]
        fill_chart_series_color(points, series_color[i])
        data_label = chart.series[i].data_labels
        data_label.show_value = True
        data_label.font.size = Pt(11)
        data_label.font.bold = True
        data_label.font.color.rgb = RGBColor(0, 0, 0)


def format_chart_2_properties(chart, y_axis, label_flag=None):
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.format.line.fill.background()
    tick_labels = chart.value_axis.tick_labels
    tick_labels.font.size = Pt(10)
    tick_labels.font.color.rgb = RGBColor(123, 123, 123)
    chart.category_axis.has_major_gridlines = False
    tick_labels = chart.category_axis.tick_labels
    tick_labels.font.size = Pt(10)
    tick_labels.font.color.rgb = RGBColor(123, 123, 123)
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_label = plot.data_labels
    data_label.font.size = Pt(11)
    data_label.font.bold = True
    data_label.position = XL_LABEL_POSITION.OUTSIDE_END
    series_color = [(191, 191, 191)]
    for i in range(len(y_axis)):
        points = plot.series[i]
        fill_chart_series_color(points, series_color[i])


def format_chart_3_properties(chart, y_axis, label_flag=None):
    chart.has_legend = False
    value_axis = chart.value_axis
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.has_minor_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH
    tick_labels = chart.value_axis.tick_labels
    tick_labels.font.size = Pt(10)
    tick_labels.font.color.rgb = RGBColor(123, 123, 123)
    chart.value_axis.format.line.fill.background()
    category_axis = chart.category_axis
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    chart.category_axis.format.line.fill.background()
    series_color = [(155, 187, 89), (247, 150, 70), (193, 81, 78)]  # Green,Orange,Red
    for i in range(len(y_axis)):
        data_label = chart.series[i].data_labels
        if label_flag:
            show_label = label_flag[i]
        else:
            show_label = True
        data_label.show_value = show_label
        data_label.font.size = Pt(11)
        data_label.font.bold = True
        data_label.font.color.rgb = RGBColor(series_color[i][0], series_color[i][1], series_color[i][2])
        if i == 0:
            data_label.position = XL_LABEL_POSITION.BELOW
        elif i == 1:
            data_label.position = XL_LABEL_POSITION.ABOVE
        chart.series[i].format.line.color.rgb = RGBColor(series_color[i][0], series_color[i][1], series_color[i][2])


def create_line_chart(slide, x_axis, y_axis, values, chart_dim, cp, title=None, label_flag=None):
    chart_data = CategoryChartData()
    chart_data.categories = x_axis
    for i in range(len(y_axis)):
        chart_data.add_series(y_axis[i], values[i])
    x, y, cx, cy = chart_dim[0], chart_dim[1], chart_dim[2], chart_dim[3]
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    if title:
        title_tf = chart.chart_title.text_frame
        title_tf.text = title
        title_font = title_tf.paragraphs[0].font
        title_font.size, title_font.bold, title_font.name = Pt(18), False, "Proxima Nova"
    cp(chart, y_axis, label_flag)


def create_bar_chart(slide, x_axis, y_axis, values, chart_dim, chart_type, cp, title=None, label_flag=None):
    chart_data = CategoryChartData()
    chart_data.categories = x_axis
    for i in range(len(y_axis)):
        chart_data.add_series(y_axis[i], values[i])
    x, y, cx, cy = chart_dim[0], chart_dim[1], chart_dim[2], chart_dim[3]
    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    if title:
        title_tf = chart.chart_title.text_frame
        title_tf.text = title
        title_font = title_tf.paragraphs[0].font
        title_font.size, title_font.bold, title_font.name = Pt(14), True, "Proxima Nova"
        title_font.color.rgb = RGBColor(128, 128, 128)
    cp(chart, y_axis, label_flag)


def create_custom_legends(slide, table_dim, legend_dct):
    shape = slide.shapes.add_table(1, 10, table_dim[0], table_dim[1], table_dim[2], table_dim[3])
    table = shape.table
    for i in range(10):
        cell = table.cell(0, i)
        try:
            table.columns[i].width = legend_dct[i]["width"]
        except:
            pass
        try:
            rgb = legend_dct[i]["cell_color"]
            fill_table_cell_color(cell, rgb)
        except:
            pass
        try:
            text = legend_dct[i]["text"]
            size = legend_dct[i]["font_size"]
            rgb = legend_dct[i]["font_color"]
            bold = legend_dct[i]["bold"]
            align = legend_dct[i]["align"]
            tf = cell.text_frame
            font = format_cell_text(cell, tf, text, size, bold, align=align)
            font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        except:
            pass
        try:
            rotation = legend_dict[i]["rotation"]
            tcpr = cell._tc.get_or_add_tcPr()
            tcpr.set('vert', 'vert270')
        except:
            pass


def calculate_time(func):
    def wrapper(*args, **kwargs):
        start_time = time.time()
        result = func(*args, **kwargs)
        print("**Time-Taken**: {} sec".format(time.time() - start_time))
        return result

    return wrapper


def create_presentation_with_branding_image(image_paths, n_slide):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(6)
    blank_slide_layout = prs.slide_layouts[6]
    for i in range(n_slide):
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = 0
        if i == 0:
            slide.shapes.add_picture(image_paths[0], left, top, width=prs.slide_width, height=prs.slide_height)
        elif i == 1:
            slide.shapes.add_picture(image_paths[1], left, top, width=prs.slide_width, height=prs.slide_height)
        else:
            slide.shapes.add_picture(image_paths[2], left, top, width=prs.slide_width, height=prs.slide_height)
    return prs


def reformat_slide_1(prs, df_fdict):
    slide = prs.slides[0]
    title = get_title(df_fdict)
    subtitle = "May 4, 2020"
    add_text_box(slide, [Cm(1.3), Cm(10), Inches(9), Inches(1)], title, 32, True)
    add_text_box(slide, [Cm(1.3), Cm(11.5), Inches(9), Inches(1)], subtitle, 28)
    return prs, title


def reformat_slide_2(prs):
    slide = prs.slides[1]
    title = "Culture & Talent"
    font = add_text_box(slide, [Cm(1.3), Cm(6.5), Inches(9), Inches(1)], title, 36, True)
    fill_text_color(font, [255, 255, 255])
    return prs


def reformat_slide_3(prs, df_fdict):
    slide = prs.slides[2]
    heading = "QoQ Key Culture & Talent Metrics"
    table_dict = get_json_data_for_table_1(["Q4 FY20", "Q1 FY21"], df_fdict)
    add_text_box(slide, [Cm(1.3), Inches(0.2), Inches(9), Inches(1)], heading, 20, True, align_center=False)
    create_table_box(slide, [Inches(3.2), Inches(1), Inches(2.4), Inches(4)], 9, table_dict, 2)
    return prs


def reformat_slide_4(prs, df_fdict):
    slide = prs.slides[3]
    heading = "QoQ Attrition Overview"
    add_text_box(slide, [Cm(1.3), Inches(0.2), Inches(9), Inches(1)], heading, 20, True, align_center=False)
    x_axis, y_axis, values, label_flags = reformat_dict_attr_for_attrition_chart(df_fdict)
    create_bar_chart(slide, x_axis, y_axis, values, [Cm(1), Inches(1), Inches(4.8), Inches(4.8)],
                     XL_CHART_TYPE.COLUMN_STACKED, format_chart_1_properties, label_flag=label_flags)
    chart_table = get_json_data_for_table_2(df_fdict)
    create_table_box(slide, [Inches(5.1), Inches(1), Inches(3.8), Inches(3.5)], 9, chart_table, 1.2)
    return prs


def reformat_slide_5(prs, ehi_db):
    slide = prs.slides[4]
    heading = "EHI Trend"
    add_text_box(slide, [Cm(1.3), Inches(0.2), Inches(9), Inches(1)], heading, 20, True, align_center=False)
    x_axis, y_axis_1, values_1, y_axis_2, values_2, data_labels = get_dict_attr_from_dataframe_for_ehi(ehi_db)
    create_bar_chart(slide, x_axis, y_axis_1, values_1, [Cm(3.3), Inches(0.7), Inches(7.2), Inches(4)],
                     XL_CHART_TYPE.COLUMN_CLUSTERED, format_chart_2_properties,
                     "Average EHI and Category Distribution over EDAP Cycles")
    create_line_chart(slide, x_axis, y_axis_2, values_2, [Cm(3.6), Inches(1.36), Inches(7.4), Inches(3.15)],
                      format_chart_3_properties)
    create_custom_legends(slide, [Inches(1.3), Inches(5), Inches(2), Inches(0.4)], legend_dict)
    add_text_box(slide, [Cm(2.8), Inches(2.5), Inches(1), Inches(1)], "Average EHI", 11, False, align_center=False,
                 rotation=-90, rgb=(123, 123, 123))
    add_text_box(slide, [Inches(8.7), Inches(2.9), Inches(1), Inches(1)], "Number of Employees", 11, False,
                 align_center=False, rotation=-90, rgb=(123, 123, 123))
    add_text_box(slide, [Inches(4.5), Inches(4.6), Inches(1), Inches(0.2)], "EHI Cycle", 11, False,
                 align_center=False, rgb=(123, 123, 123))
    return prs


@calculate_time
def create_ops_review_presentation(images, df_fdict, ehi_db, ppt_path=None):
    prs = create_presentation_with_branding_image(images, 5)
    prs, title = reformat_slide_1(prs, df_fdict)
    prs = reformat_slide_2(prs)
    prs = reformat_slide_3(prs, df_fdict)
    prs = reformat_slide_4(prs, df_fdict)
    prs = reformat_slide_5(prs, ehi_db)
    file_name = title + " - Ops Review.pptx"
    if ppt_path:
        prs.save(os.path.join(ppt_path, file_name))
        print("{} created".format(file_name))
    return None


parent_dir = Path(os.getcwd()).parent
deck_path = os.path.join(parent_dir, "Decks")
files_path = os.path.join(parent_dir, "Files")
images_path = os.path.join(parent_dir, "Images")
image_names = [os.path.join(images_path, "title_page.png"), os.path.join(images_path, "branding_page.png"),
               os.path.join(images_path, "content_page.png")]

db_dict = get_dataframe_dict(files_path, ["Q1 FY20", "Q2 FY20", "Q3 FY20", "Q4 FY20", "Q1 FY21"])

uid = "53HWTWIBW"
db_fdict = filter_dataframe_dict_on_uid(uid, db_dict)
db = read_ehi_data(os.path.join(files_path, "ehi_example_data_fozia.xlsx"))
create_ops_review_presentation(image_names, db_fdict, db, deck_path)

uid = "LA0UT4HLD"
db_fdict = filter_dataframe_dict_on_uid(uid, db_dict)
db = read_ehi_data(os.path.join(files_path, "ehi_example_data_hadrian.xlsx"))
create_ops_review_presentation(image_names, db_fdict, db, deck_path)

uid = "ZJG9NQVX8"
db_fdict = filter_dataframe_dict_on_uid(uid, db_dict)
db = read_ehi_data(os.path.join(files_path, "ehi_example_data_heidi.xlsx"))
create_ops_review_presentation(image_names, db_fdict, db, deck_path)
